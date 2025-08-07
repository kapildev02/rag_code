from app.db.mongodb import (
    connect_to_mongodb,
    document_collection,
    get_fs
)
from app.core.rabbitmq_client import rabbitmq_client
from app.core.config import settings
import asyncio
import aio_pika
import json
from bson import ObjectId
import markdownify
import os
import fitz 
import io
import logging
from docx import Document
from pptx import Presentation
import xml.etree.ElementTree as ET

async def on_message(message: aio_pika.IncomingMessage):
    try:
        decoded_msg = json.loads(message.body.decode())
        
        # doc id
        doc_id = decoded_msg["doc_id"]
        
        # check if doc id present or not
        if not doc_id:
            return
        
        # fetch document details
        doc = await document_collection().find_one({
            "_id": ObjectId(doc_id)
        })
        
        # check if the document present
        if not doc:
            return
        
        # get gridfs raw file id
        raw_gridfs_id = doc["raw_gridfs_id"]
        
        # check raw_gridfs_id present
        if not raw_gridfs_id:
            return
        
        # fetch the file from gridfs
        raw_gridfs_file = await get_fs().open_download_stream(raw_gridfs_id)
        raw_gridfs_file_bytes = await raw_gridfs_file.read()

         # Determine file type and extract text
        filename = doc['filename']
        file_extension = filename.rsplit('.', 1)[-1].lower()
        raw_txt = ""

        if file_extension == "pdf":
            # Open PDF using PyMuPDF
            pdf_stream = io.BytesIO(raw_gridfs_file_bytes)
            doc_pdf = fitz.open(stream=pdf_stream, filetype="pdf")
        
        # Extract text from PDF
        # raw_txt = ""
            for page in doc_pdf:
              text = page.get_text("text")  # You can also try "html"
              raw_txt += text + "\n\n"
            
            doc_pdf.close()

         
        elif file_extension == "docx":
            # Open DOCX using python-docx
            docx_stream = io.BytesIO(raw_gridfs_file_bytes)
            document = Document(docx_stream)
            for para in document.paragraphs:
                raw_txt += para.text + "\n\n"

        elif file_extension == "pptx":
            # Open PPTX using python-pptx
            pptx_stream = io.BytesIO(raw_gridfs_file_bytes)
            presentation = Presentation(pptx_stream)
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        raw_txt += shape.text + "\n\n"

        elif file_extension == "xml":
            # Parse XML using ElementTree
            xml_stream = io.BytesIO(raw_gridfs_file_bytes)
            tree = ET.parse(xml_stream)
            root = tree.getroot()
            raw_txt = ET.tostring(root, encoding='unicode', method='text')

        elif file_extension == "txt":
            # Directly use the text content
            raw_txt = raw_gridfs_file_bytes.decode('utf-8')

        # convert to markedown
        md_text = markdownify.markdownify(raw_txt)
        md_file_name = f"{doc['filename']}.md"
        
        
        # Save Markdown file locally
        local_dir = os.path.join("uploaded_local_files")
        os.makedirs(local_dir, exist_ok=True)
        local_file_path = os.path.join(local_dir, md_file_name)

        with open(local_file_path, "w", encoding="utf-8") as f:
            f.write(md_text)
            
        print(f"md_file_name: {md_file_name}")
        
        # Manually ack after successful processing
        # await message.ack()
    except Exception as e:
        print(f"Processing failed: {e}")
        # You can nack and optionally requeue the message
        # await message.nack(requeue=True)


async def main():
    await connect_to_mongodb()
    await rabbitmq_client.connect()
    await rabbitmq_client.channel.set_qos(prefetch_count=1)
    
    # Start consuming (non-blocking, async)
    await rabbitmq_client.consume_message(settings.FILE_PROCESSING_CHANNEL, on_message)
    await asyncio.Future()
    
    
if __name__ == "__main__":
    asyncio.run(main())